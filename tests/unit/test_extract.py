"""텍스트 추출기 단위테스트 (FR-006 `.txt`/`.md`, FR-007 `.docx`, v0.2 U1 `.pdf`)."""

from __future__ import annotations

import os
from datetime import date, datetime
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Emu
from pypdf import PdfReader, PdfWriter

from corpbrain.core.config import SUPPORTED_EXTENSIONS
from corpbrain.core.extract import (
    EXTRACTORS,
    ExtractionError,
    _open_failure_detail,
    extract_text,
    prepare_summary_input,
)
from corpbrain.core.models import SkipReason

MAX_CHARS = 12000

#: 완료의 정의(스펙 §3)가 요구하는 픽스처 코퍼스 — 통합테스트(FR-018)도 같은 폴더를 쓴다.
FIXTURE_CORPUS = Path(__file__).resolve().parents[1] / "fixtures" / "sample_corpus"

#: 권한 거부 재현은 POSIX 권한 비트가 적용되는 비-root 환경에서만 가능하다. `os.geteuid`는
#: POSIX 전용이므로 단락 평가로 Windows 수집 크래시를 피한다 (test_scanner.py와 동일 관용구).
needs_posix_permissions = pytest.mark.skipif(
    os.name != "posix" or os.geteuid() == 0,
    reason="POSIX 권한 비트가 적용되는 비-root 환경에서만 권한 거부를 재현할 수 있다",
)


def _write_docx(path: Path, paragraphs: list[str]) -> None:
    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    document.save(str(path))


def _write_text_pdf(path: Path, text: str) -> None:
    """추출 가능한 텍스트 레이어가 있는 최소 단일 페이지 PDF를 만든다 (외부 의존성 없음).

    한글은 base14 폰트로 추출되지 않으므로 텍스트는 ASCII만 사용한다. 바이트 오프셋을 계산해
    올바른 xref를 써서 pypdf의 xref 재구성 경고 없이 읽히게 한다.
    """
    stream = b"BT /F1 24 Tf 72 700 Td (" + text.encode("ascii") + b") Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
        ),
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets: list[int] = []
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += str(index).encode() + b" 0 obj\n" + body + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 " + str(len(objects) + 1).encode() + b"\n0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += b"trailer\n<< /Size " + str(len(objects) + 1).encode() + b" /Root 1 0 R >>\n"
    out += b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF\n"
    path.write_bytes(bytes(out))


def _write_encrypted_pdf(path: Path, source_pdf: Path) -> None:
    """`source_pdf`를 사용자 암호로 암호화해 저장한다 (추출 시 암호화로 스킵되어야 한다)."""
    writer = PdfWriter()
    writer.append(PdfReader(str(source_pdf)))
    writer.encrypt("secret")
    with path.open("wb") as stream:
        writer.write(stream)


def _write_blank_pdf(path: Path) -> None:
    """텍스트 레이어가 없는 PDF(스캔 이미지 PDF 근사) — 추출 결과가 빈 문자열이어야 한다."""
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with path.open("wb") as stream:
        writer.write(stream)


#: OLE 복합문서(CFBF) 매직 넘버 — 암호화된 OOXML 파일과 구형 `.xls`/`.ppt`가 공유한다.
OLE_SIGNATURE_BYTES = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def test_ole_signature_yields_encrypted_or_legacy_detail(tmp_path: Path) -> None:
    """OLE 시그니처로 시작하는 파일은 「암호화되었거나 구형 이진 포맷」 detail을 받는다 (스펙 §4.3.1).

    암호화된 `.xlsx`/`.pptx`는 zip이 아니라 OLE 복합문서로 감싸져 라이브러리가 「zip을 열지
    못했다」는 예외만 던진다 — 손상된 파일과 예외 종류가 같아 구분되지 않는다. 판정은 예외
    메시지 문자열이 아니라 파일 시그니처로 한다.
    """
    source = tmp_path / "locked.xlsx"
    source.write_bytes(OLE_SIGNATURE_BYTES + b"\x00" * 32)

    detail = _open_failure_detail(source)

    assert "암호화되었거나 구형 이진 포맷" in detail
    assert str(source) in detail


def test_non_ole_file_yields_plain_open_failure_detail(tmp_path: Path) -> None:
    """OLE 시그니처가 아니면 손상·파싱 실패 문구를 그대로 쓴다 — 두 원인이 다른 detail을 받는다."""
    source = tmp_path / "broken.xlsx"
    source.write_bytes(b"PK\x03\x04not-a-real-zip")

    detail = _open_failure_detail(source)

    assert "암호화되었거나 구형 이진 포맷" not in detail
    assert str(source) in detail


def test_ole_detail_names_the_actual_extension(tmp_path: Path) -> None:
    """문구의 포맷 이름은 실제 확장자에서 온다 — `.xlsm`이 「xlsx」로 보고되지 않는다."""
    source = tmp_path / "macro.xlsm"
    source.write_bytes(OLE_SIGNATURE_BYTES)

    assert "xlsm" in _open_failure_detail(source)


def test_ole_probe_survives_an_unreadable_file(tmp_path: Path) -> None:
    """시그니처를 읽지 못해도(파일 소실·권한) 예외 없이 기본 문구로 떨어진다.

    이 헬퍼는 **이미 실패한 경로**에서만 불린다. 여기서 새 예외를 올리면 원래의 추출 실패가
    다른 예외로 뒤덮여 스킵 사유가 바뀐다.
    """
    missing = tmp_path / "gone.xlsx"

    detail = _open_failure_detail(missing)

    assert "암호화되었거나 구형 이진 포맷" not in detail


def test_dispatch_table_keys_match_supported_extensions() -> None:
    """확장자 디스패치 매핑의 키 집합이 `SUPPORTED_EXTENSIONS`와 정확히 같다 (스펙 §3 항목7).

    지원 포맷 목록이 `config.SUPPORTED_EXTENSIONS`와 `extract.py`에 따로 정의돼 있어 둘이
    어긋나도 아무도 검증하지 않던 구멍을 닫는다 — 지금까지의 유일한 방어는 `test_scanner.py`의
    리터럴 단언뿐이었고, 그것은 상수 쪽만 본다. 한쪽에만 확장자를 더하면 스캐너가 통과시킨
    파일을 추출기가 「지원하지 않는 확장자」로 되던지거나(추출 실패로 위장된 미지원), 추출기는
    아는데 스캐너가 걸러 영원히 불리지 않는 죽은 코드가 된다.
    """
    assert set(EXTRACTORS) == set(SUPPORTED_EXTENSIONS)


def test_dispatch_table_maps_every_extension_to_a_callable() -> None:
    """매핑 값은 모두 호출 가능해야 한다 — 키만 맞고 값이 비면 위 단언이 공허해진다."""
    assert all(callable(extractor) for extractor in EXTRACTORS.values())


def _write_xlsx(path: Path, sheets: dict[str, list[list[object]]], *, hidden: tuple[str, ...] = ()) -> None:
    """시트 이름 → 행 목록으로 `.xlsx`/`.xlsm`을 인라인 생성한다 (스펙 §3 픽스처 구성).

    픽스처 생성 헬퍼는 `tests/conftest.py`나 공유 모듈로 빼지 않고 단위·통합 테스트 파일에
    각각 둔다 — `.pdf`가 이미 같은 선택을 했고(위 `_write_text_pdf` 참조), 공유 인프라를 새로
    들이면 `.pdf`는 복제, 오피스는 공유가 되어 관용구가 둘로 갈린다 (스펙 §3 「헬퍼 배치」).
    """
    workbook = Workbook()
    workbook.remove(workbook.active)
    for title, rows in sheets.items():
        worksheet = workbook.create_sheet(title)
        for row in rows:
            worksheet.append(row)
        if title in hidden:
            worksheet.sheet_state = "hidden"
    workbook.save(str(path))



def _blank_slide(presentation: Presentation) -> object:
    """도형이 하나도 없는 빈 레이아웃 슬라이드를 더한다 (레이아웃 6 = Blank)."""
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _add_textbox(slide: object, text: str, *, top: int = 0) -> object:
    box = slide.shapes.add_textbox(Emu(0), Emu(top), Emu(1_000_000), Emu(500_000))
    box.text_frame.text = text
    return box



def test_txt_extraction_returns_plain_text(tmp_path: Path) -> None:
    source = tmp_path / "note.txt"
    source.write_text("첫 줄\n둘째 줄\n", encoding="utf-8")

    assert extract_text(source, MAX_CHARS) == "첫 줄\n둘째 줄\n"


def test_md_is_extracted_verbatim_without_markdown_parsing(tmp_path: Path) -> None:
    source = tmp_path / "doc.md"
    source.write_text("# 제목\n\n- 항목 1\n", encoding="utf-8")

    assert extract_text(source, MAX_CHARS) == "# 제목\n\n- 항목 1\n"


def test_plaintext_reads_only_leading_max_chars(tmp_path: Path) -> None:
    source = tmp_path / "long.txt"
    source.write_text("가" * 20_000, encoding="utf-8")

    extracted = extract_text(source, MAX_CHARS)

    assert len(extracted) == MAX_CHARS
    assert extracted == "가" * MAX_CHARS


def test_empty_file_yields_empty_string(tmp_path: Path) -> None:
    source = tmp_path / "empty.txt"
    source.write_text("", encoding="utf-8")

    assert extract_text(source, MAX_CHARS) == ""


def test_broken_encoding_does_not_crash(tmp_path: Path) -> None:
    source = tmp_path / "broken.txt"
    source.write_bytes(b"\xff\xfe valid tail")

    extracted = extract_text(source, MAX_CHARS)

    assert "valid tail" in extracted


def test_unsupported_extension_raises_extraction_error(tmp_path: Path) -> None:
    # `.xlsx`는 v0.8에서 지원 포맷이 되었다. 구형 `.xls`(BIFF)는 여전히 비목표다 (v0.8 §2).
    source = tmp_path / "sheet.xls"
    source.write_bytes(b"PK\x03\x04")

    with pytest.raises(ExtractionError):
        extract_text(source, MAX_CHARS)


@needs_posix_permissions
def test_permission_denied_is_reported_as_extraction_error(tmp_path: Path) -> None:
    source = tmp_path / "secret.txt"
    source.write_text("내용", encoding="utf-8")
    source.chmod(0o000)

    try:
        with pytest.raises(ExtractionError) as excinfo:
            extract_text(source, MAX_CHARS)
    finally:
        source.chmod(0o600)

    assert isinstance(excinfo.value.__cause__, PermissionError)


def test_docx_paragraphs_are_joined_in_order(tmp_path: Path) -> None:
    source = tmp_path / "report.docx"
    _write_docx(source, ["첫 문단", "둘째 문단", "셋째 문단"])

    assert extract_text(source, MAX_CHARS) == "첫 문단\n둘째 문단\n셋째 문단"


def test_docx_extraction_stops_at_max_chars(tmp_path: Path) -> None:
    source = tmp_path / "big.docx"
    _write_docx(source, ["나" * 500 for _ in range(20)])

    extracted = extract_text(source, 1000)

    assert len(extracted) == 1000


def test_corrupted_docx_raises_extraction_error(tmp_path: Path) -> None:
    source = tmp_path / "corrupt.docx"
    source.write_bytes(b"not a real docx package")

    with pytest.raises(ExtractionError):
        extract_text(source, MAX_CHARS)


def test_empty_docx_yields_empty_string(tmp_path: Path) -> None:
    source = tmp_path / "blank.docx"
    _write_docx(source, [])

    assert extract_text(source, MAX_CHARS) == ""


# --- v0.2 U1: .pdf 텍스트 레이어 추출 (스펙 §4.1·§5) -----------------------------


def test_pdf_text_layer_is_extracted(tmp_path: Path) -> None:
    source = tmp_path / "report.pdf"
    _write_text_pdf(source, "CorpBrain PDF text layer sample.")

    assert "CorpBrain PDF text layer sample." in extract_text(source, MAX_CHARS)


def test_pdf_extraction_stops_at_max_chars(tmp_path: Path) -> None:
    source = tmp_path / "long.pdf"
    _write_text_pdf(source, "CorpBrain " * 60)  # 추출 결과가 50자보다 충분히 길다

    assert len(extract_text(source, 50)) == 50


def test_encrypted_pdf_raises_extraction_error(tmp_path: Path) -> None:
    plain = tmp_path / "plain.pdf"
    _write_text_pdf(plain, "secret contents")
    locked = tmp_path / "locked.pdf"
    _write_encrypted_pdf(locked, plain)

    with pytest.raises(ExtractionError, match="암호화"):
        extract_text(locked, MAX_CHARS)


def test_encrypted_pdf_is_skipped_as_extraction_failed(tmp_path: Path) -> None:
    plain = tmp_path / "plain.pdf"
    _write_text_pdf(plain, "secret contents")
    locked = tmp_path / "locked.pdf"
    _write_encrypted_pdf(locked, plain)

    prepared = prepare_summary_input(locked, MAX_CHARS)

    assert prepared.text is None
    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.EXTRACTION_FAILED
    assert "암호화" in prepared.skipped.detail


def test_pdf_without_text_layer_yields_empty_string(tmp_path: Path) -> None:
    source = tmp_path / "scanned.pdf"
    _write_blank_pdf(source)

    assert extract_text(source, MAX_CHARS) == ""


def test_pdf_without_text_layer_is_skipped_as_empty(tmp_path: Path) -> None:
    source = tmp_path / "scanned.pdf"
    _write_blank_pdf(source)

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.text is None
    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.EMPTY_DOCUMENT


def test_corrupted_pdf_is_skipped_as_extraction_failed(tmp_path: Path) -> None:
    source = tmp_path / "corrupt.pdf"
    source.write_bytes(b"%PDF-1.4 this is not a real pdf body")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.text is None
    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.EXTRACTION_FAILED
    assert prepared.skipped.detail


# --- FR-008: 길이 제한 + 빈문서·추출실패 스킵 판정 -------------------------------


def test_oversized_document_is_truncated_not_skipped(tmp_path: Path) -> None:
    source = tmp_path / "huge.txt"
    source.write_text("다" * 20_000, encoding="utf-8")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.skipped is None
    assert prepared.text is not None
    assert len(prepared.text) == MAX_CHARS


def test_whitespace_only_document_is_skipped_as_empty(tmp_path: Path) -> None:
    source = tmp_path / "blank.md"
    source.write_text("   \n\t\n", encoding="utf-8")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.text is None
    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.EMPTY_DOCUMENT


def test_zero_byte_document_is_skipped_as_empty(tmp_path: Path) -> None:
    source = tmp_path / "zero.txt"
    source.write_bytes(b"")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.EMPTY_DOCUMENT


def test_corrupted_docx_is_skipped_as_extraction_failed(tmp_path: Path) -> None:
    source = tmp_path / "corrupt.docx"
    source.write_bytes(b"not a real docx package")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.text is None
    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.EXTRACTION_FAILED
    assert prepared.skipped.detail


@needs_posix_permissions
def test_unreadable_document_is_skipped_as_permission_denied(tmp_path: Path) -> None:
    source = tmp_path / "locked.txt"
    source.write_text("내용", encoding="utf-8")
    source.chmod(0o000)

    try:
        prepared = prepare_summary_input(source, MAX_CHARS)
    finally:
        source.chmod(0o600)

    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.PERMISSION_DENIED


def test_normal_document_passes_through_with_text(tmp_path: Path) -> None:
    source = tmp_path / "ok.txt"
    source.write_text("정상 문서 본문", encoding="utf-8")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.skipped is None
    assert prepared.text == "정상 문서 본문"


# --- FR-017: 커밋된 픽스처 코퍼스 기반 검증 -------------------------------------


def test_fixture_txt_and_md_extract_plain_text() -> None:
    assert "샘플 텍스트 문서입니다." in extract_text(FIXTURE_CORPUS / "normal.txt", MAX_CHARS)
    assert "# 마크다운 제목" in extract_text(FIXTURE_CORPUS / "guide.md", MAX_CHARS)


def test_fixture_docx_extracts_paragraphs_in_order() -> None:
    extracted = extract_text(FIXTURE_CORPUS / "sub" / "report.docx", MAX_CHARS)

    assert extracted.splitlines()[0] == "분기 실적 보고"
    assert "신규 고객이 늘었다." in extracted


def test_fixture_oversized_document_is_truncated_to_max_chars() -> None:
    prepared = prepare_summary_input(FIXTURE_CORPUS / "oversized.txt", MAX_CHARS)

    assert prepared.skipped is None
    assert prepared.text is not None
    assert len(prepared.text) == MAX_CHARS


def test_fixture_empty_document_is_skipped() -> None:
    prepared = prepare_summary_input(FIXTURE_CORPUS / "empty.txt", MAX_CHARS)

    assert prepared.skipped is not None
    assert prepared.skipped.reason == SkipReason.EMPTY_DOCUMENT


def test_fixture_unsupported_extension_is_rejected() -> None:
    with pytest.raises(ExtractionError):
        extract_text(FIXTURE_CORPUS / "photo.jpg", MAX_CHARS)


# --- 엑셀 `.xlsx`/`.xlsm` (v0.8 §4.2) --------------------------------------------


def test_xlsx_sheets_are_bounded_and_rows_are_tab_separated(tmp_path: Path) -> None:
    """내용이 있는 시트마다 경계 줄이 정확히 1회 나오고, 셀은 탭으로 구분된다 (§3 항목2).

    경계가 없으면 서로 다른 표의 행이 이어 붙어 한 표로 읽힌다. 시트명 자체도 요약에 쓸모
    있는 신호다.
    """
    source = tmp_path / "book.xlsx"
    _write_xlsx(
        source,
        {
            "인사": [["이름", "부서"], ["김철수", "인사팀"]],
            "예산": [["항목", "금액"]],
        },
    )

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == "[시트: 인사]\n이름\t부서\n김철수\t인사팀\n[시트: 예산]\n항목\t금액"
    assert extracted.count("[시트: 인사]") == 1


def test_xlsx_sheet_without_content_gets_no_boundary_line(tmp_path: Path) -> None:
    """값이 없는 시트에는 경계 줄도 내지 않는다 (스펙 §4.2 T1).

    경계 줄은 「뒤에 내용이 온다」는 신호이므로 뒤가 비면 낼 것이 없다. 항상 내면 값이 전부
    비는 파일의 추출 결과가 경계 줄만 남아 호출자의 `text.strip()` 검사를 통과해 버린다.
    """
    source = tmp_path / "book.xlsx"
    _write_xlsx(source, {"빈시트": [[None, None]], "내용": [["값"]]})

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == "[시트: 내용]\n값"
    assert "빈시트" not in extracted


def test_xlsx_hidden_sheet_content_is_excluded(tmp_path: Path) -> None:
    """숨긴 시트는 내용도 경계 줄도 나오지 않는다 — `sheet_state`는 읽기 전용 개봉에서도 판정된다."""
    source = tmp_path / "book.xlsx"
    _write_xlsx(source, {"보이는": [["공개"]], "숨긴": [["기밀"]]}, hidden=("숨긴",))

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == "[시트: 보이는]\n공개"
    assert "기밀" not in extracted


def test_xlsx_hidden_row_content_is_included(tmp_path: Path) -> None:
    """숨긴 **행**은 제외되지 않는다 — `read_only` 개봉에서 판별할 수 없기 때문이다.

    `openpyxl` 3.1.5가 돌려주는 `ReadOnlyWorksheet`에는 `row_dimensions` 속성이 아예 없다
    (2026-08-27 U1 확인). 스펙 §4.2 T3의 조건부 표가 「얻을 수 없다」 행으로 확정된 결과이며,
    `read_only=False`로 바꾸는 선택지는 「파일을 통째로 메모리에 올리지 않는다」와 정면으로
    부딪쳐 택하지 않는다. 알려진 한계를 테스트로 고정해 둔다.
    """
    source = tmp_path / "book.xlsx"
    _write_xlsx(source, {"시트": [["보임1"], ["숨김"], ["보임2"]]})

    workbook = Workbook()  # 위에서 만든 파일에 행 숨김만 덧입힌다
    del workbook
    from openpyxl import load_workbook

    editable = load_workbook(str(source))
    editable["시트"].row_dimensions[2].hidden = True
    editable.save(str(source))

    extracted = extract_text(source, MAX_CHARS)

    assert "숨김" in extracted


def test_xlsx_cell_values_are_normalised(tmp_path: Path) -> None:
    """셀 값의 문자열화 4규칙 (스펙 §4.2 T2 표).

    `openpyxl`이 돌려주는 것은 문자열이 아니라 파이썬 객체다. `120000.0`·
    `2026-03-31 00:00:00` 같은 표기는 LLM이 요약문에 그대로 옮겨 적는 종류의 잡음이다.
    """
    source = tmp_path / "book.xlsx"
    _write_xlsx(
        source,
        {
            # 엑셀 셀의 날짜는 tz 정보가 없다 — `openpyxl`은 tz-aware `datetime` 저장을 거부하므로
            # 여기서 naive 를 쓰는 것이 실제 셀 값의 모습이다 (DTZ001 을 그래서 끈다).
            "값": [
                [date(2026, 3, 31), datetime(2026, 3, 31, 14, 30)],  # noqa: DTZ001
                [datetime(2026, 3, 31, 0, 0), 120000.0],  # noqa: DTZ001
                [1.5, True],
            ]
        },
    )

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == (
        "[시트: 값]\n"
        "2026-03-31\t2026-03-31 14:30\n"
        "2026-03-31\t120000\n"
        "1.5\tTrue"
    )


def test_xlsx_blank_rows_are_skipped_and_trailing_empty_cells_leave_no_tab(
    tmp_path: Path,
) -> None:
    """전 셀이 빈 행은 건너뛰고, 행 끝의 빈 셀은 탭을 남기지 않는다 (스펙 §4.2 T2).

    앞쪽·중간의 빈 셀은 열 위치 정보이므로 탭으로 남긴다.
    """
    source = tmp_path / "book.xlsx"
    _write_xlsx(
        source,
        {"시트": [["a", "b", None], [None, None, None], [None, "c", None], ["d"]]},
    )

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == "[시트: 시트]\na\tb\n\tc\nd"


def test_xlsx_formula_without_cached_value_yields_empty_string(tmp_path: Path) -> None:
    """수식만 있고 계산값 캐시가 없는 통합문서는 빈 문자열이 된다 (스펙 §3 항목3·§5).

    Excel이 저장하지 않은 `.xlsx`(스크립트·LibreOffice 생성)에는 `data_only=True`가 읽을
    캐시가 없다. 값을 수식 문자열로 대체하려고 파일을 두 번 열지 않는다.
    """
    source = tmp_path / "formula.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "계산"
    worksheet["A1"] = "=SUM(B1:B9)"
    worksheet["A2"] = "=A1*2"
    workbook.save(str(source))

    assert extract_text(source, MAX_CHARS) == ""


def test_xlsx_formula_only_workbook_is_skipped_as_empty(tmp_path: Path) -> None:
    """그 빈 문자열이 호출자에게서 `empty_document`가 된다 — 빈 판정 책임은 호출자에 남는다."""
    source = tmp_path / "formula.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "=SUM(B1:B9)"
    workbook.save(str(source))

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.text is None
    assert prepared.skipped is not None
    assert prepared.skipped.reason is SkipReason.EMPTY_DOCUMENT


def test_xlsx_extraction_stops_at_max_chars(tmp_path: Path) -> None:
    """상한에 닿으면 즉시 멈춘다 — 통합문서를 끝까지 훑지 않는다 (§3 항목9)."""
    source = tmp_path / "big.xlsx"
    _write_xlsx(source, {"시트": [["가" * 200] for _ in range(50)]})

    extracted = extract_text(source, 500)

    assert len(extracted) <= 500


def test_xlsm_is_extracted_by_the_same_extractor(tmp_path: Path) -> None:
    """`.xlsm`은 `openpyxl`이 공식 지원하므로 확장자 한 줄 추가로 끝난다 (스펙 §4.1)."""
    source = tmp_path / "macro.xlsm"
    _write_xlsx(source, {"시트": [["매크로 통합문서"]]})

    assert extract_text(source, MAX_CHARS) == "[시트: 시트]\n매크로 통합문서"


def test_ole_xlsx_and_corrupted_xlsx_get_different_details(tmp_path: Path) -> None:
    """암호화(OLE)와 손상이 **서로 다른 detail**을 받는다 (스펙 §3 항목5·§4.3.1)."""
    encrypted = tmp_path / "locked.xlsx"
    encrypted.write_bytes(OLE_SIGNATURE_BYTES + b"\x00" * 64)
    corrupted = tmp_path / "broken.xlsx"
    corrupted.write_bytes(b"PK\x03\x04 not really a zip")

    with pytest.raises(ExtractionError) as encrypted_error:
        extract_text(encrypted, MAX_CHARS)
    with pytest.raises(ExtractionError) as corrupted_error:
        extract_text(corrupted, MAX_CHARS)

    assert "암호화되었거나 구형 이진 포맷" in str(encrypted_error.value)
    assert "암호화되었거나 구형 이진 포맷" not in str(corrupted_error.value)
    assert str(encrypted_error.value) != str(corrupted_error.value)


def test_broken_xlsx_is_skipped_as_extraction_failed(tmp_path: Path) -> None:
    """손상된 통합문서는 새 스킵 사유가 아니라 기존 `extraction_failed`로 흡수된다 (스펙 §4.3)."""
    source = tmp_path / "broken.xlsx"
    source.write_bytes(b"PK\x03\x04 not really a zip")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.skipped is not None
    assert prepared.skipped.reason is SkipReason.EXTRACTION_FAILED
    assert prepared.skipped.detail


# --- PPTX (v0.8 §4.2) -------------------------------------------------------------


def test_pptx_slides_are_bounded_and_shapes_keep_file_order(tmp_path: Path) -> None:
    """내용이 있는 슬라이드마다 `[슬라이드 N]` 경계 줄이 나오고, 도형은 파일 순서 그대로다.

    좌표 기반 시각 순서 재정렬을 하지 않는다 — 같은 높이의 도형에서 순서가 흔들려 tie-break를
    또 정해야 하고 결정성이 약해진다. 아래 픽스처는 **나중에 추가한 도형을 위쪽에 두어** 좌표
    순서와 파일 순서가 어긋나게 만든다.
    """
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = _blank_slide(presentation)
    _add_textbox(slide, "먼저 추가한 아래쪽 도형", top=5_000_000)
    _add_textbox(slide, "나중에 추가한 위쪽 도형", top=0)
    second = _blank_slide(presentation)
    _add_textbox(second, "둘째 슬라이드")
    presentation.save(str(source))

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == (
        "[슬라이드 1]\n먼저 추가한 아래쪽 도형\n나중에 추가한 위쪽 도형\n"
        "[슬라이드 2]\n둘째 슬라이드"
    )


def test_pptx_empty_slide_gets_no_boundary_and_numbering_is_not_reset(
    tmp_path: Path,
) -> None:
    """이미지만 있는 슬라이드는 경계 줄도 내지 않고, 남은 슬라이드는 원본 번호를 유지한다.

    건너뛴 슬라이드가 있어도 번호를 다시 매기지 않는다 — 사용자가 파워포인트에서 보는 번호와
    어긋나면 「몇 번째 슬라이드에 있던 내용인가」를 짚을 수 없다.
    """
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    _blank_slide(presentation)  # 1번: 도형 없음
    third_carrier = _blank_slide(presentation)  # 2번: 도형 없음
    del third_carrier
    slide = _blank_slide(presentation)  # 3번: 내용 있음
    _add_textbox(slide, "세 번째 슬라이드의 내용")
    presentation.save(str(source))

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == "[슬라이드 3]\n세 번째 슬라이드의 내용"


def test_pptx_empty_deck_is_skipped_as_empty(tmp_path: Path) -> None:
    """전 슬라이드가 비면 추출 결과가 빈 문자열이 되어 호출자가 `empty_document`로 스킵한다."""
    source = tmp_path / "images_only.pptx"
    presentation = Presentation()
    _blank_slide(presentation)
    _blank_slide(presentation)
    presentation.save(str(source))

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.skipped is not None
    assert prepared.skipped.reason is SkipReason.EMPTY_DOCUMENT


def test_pptx_table_cells_are_extracted(tmp_path: Path) -> None:
    """표는 `text_frame`이 없어 따로 훑지 않으면 통째로 누락된다 — 행은 탭으로 잇는다."""
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = _blank_slide(presentation)
    frame = slide.shapes.add_table(2, 2, Emu(0), Emu(0), Emu(2_000_000), Emu(1_000_000))
    table = frame.table
    table.cell(0, 0).text = "항목"
    table.cell(0, 1).text = "금액"
    table.cell(1, 0).text = "인건비"
    table.cell(1, 1).text = "1200"
    presentation.save(str(source))

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == "[슬라이드 1]\n항목\t금액\n인건비\t1200"


def test_pptx_group_shape_text_is_extracted_recursively(tmp_path: Path) -> None:
    """그룹 도형 내부 텍스트도 재귀로 훑는다 — 그러지 않으면 그룹째 누락된다."""
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = _blank_slide(presentation)
    first = _add_textbox(slide, "그룹 안 첫째")
    second = _add_textbox(slide, "그룹 안 둘째", top=1_000_000)
    slide.shapes.add_group_shape([first, second])
    presentation.save(str(source))

    extracted = extract_text(source, MAX_CHARS)

    assert "그룹 안 첫째" in extracted
    assert "그룹 안 둘째" in extracted


def test_pptx_speaker_notes_follow_a_notes_marker(tmp_path: Path) -> None:
    """발표자 노트를 `[노트]` 줄 뒤에 덧붙인다 — 본문이 키워드뿐이고 내용이 노트에 있는 경우가 흔하다."""
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = _blank_slide(presentation)
    _add_textbox(slide, "키워드")
    slide.notes_slide.notes_text_frame.text = "실제 설명은 여기 있다"
    presentation.save(str(source))

    extracted = extract_text(source, MAX_CHARS)

    assert extracted == "[슬라이드 1]\n키워드\n[노트]\n실제 설명은 여기 있다"


def test_pptx_blank_notes_produce_no_notes_marker(tmp_path: Path) -> None:
    """노트가 있으나 공백뿐이면 `[노트]` 줄도 내지 않는다 — 경계 줄 규칙과 같다."""
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = _blank_slide(presentation)
    _add_textbox(slide, "본문")
    slide.notes_slide.notes_text_frame.text = "   \n  "
    presentation.save(str(source))

    assert extract_text(source, MAX_CHARS) == "[슬라이드 1]\n본문"


def test_pptx_extraction_does_not_create_notes_slides(tmp_path: Path) -> None:
    """노트 접근을 `has_notes_slide`로 가드한다 (스펙 §4.2 T7).

    `slide.notes_slide`에 바로 접근하면 노트 슬라이드가 없는 슬라이드에 객체가 **새로
    만들어진다**. 추출기는 전부 읽기 전용이며 그 성질이 코드에서 보여야 한다. 추출 뒤에
    다시 열어 노트 슬라이드가 생기지 않았음을 확인한다.
    """
    source = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = _blank_slide(presentation)
    _add_textbox(slide, "노트 없는 슬라이드")
    presentation.save(str(source))

    extract_text(source, MAX_CHARS)

    assert Presentation(str(source)).slides[0].has_notes_slide is False


def test_pptx_extraction_stops_at_max_chars(tmp_path: Path) -> None:
    """상한에 닿으면 즉시 멈춘다 — 프레젠테이션을 끝까지 훑지 않는다 (§3 항목9)."""
    source = tmp_path / "big.pptx"
    presentation = Presentation()
    for _ in range(20):
        slide = _blank_slide(presentation)
        _add_textbox(slide, "가" * 200)
    presentation.save(str(source))

    assert len(extract_text(source, 500)) <= 500


def test_ole_pptx_and_corrupted_pptx_get_different_details(tmp_path: Path) -> None:
    """암호화(OLE)와 손상이 서로 다른 detail을 받는다 — 엑셀과 같은 판정을 공유한다 (§4.3.1)."""
    encrypted = tmp_path / "locked.pptx"
    encrypted.write_bytes(OLE_SIGNATURE_BYTES + b"\x00" * 64)
    corrupted = tmp_path / "broken.pptx"
    corrupted.write_bytes(b"PK\x03\x04 not really a zip")

    with pytest.raises(ExtractionError) as encrypted_error:
        extract_text(encrypted, MAX_CHARS)
    with pytest.raises(ExtractionError) as corrupted_error:
        extract_text(corrupted, MAX_CHARS)

    assert "암호화되었거나 구형 이진 포맷" in str(encrypted_error.value)
    assert "암호화되었거나 구형 이진 포맷" not in str(corrupted_error.value)


def test_broken_pptx_is_skipped_as_extraction_failed(tmp_path: Path) -> None:
    """손상된 프레젠테이션도 기존 `extraction_failed`로 흡수된다 — 새 스킵 사유를 만들지 않는다."""
    source = tmp_path / "broken.pptx"
    source.write_bytes(b"PK\x03\x04 not really a zip")

    prepared = prepare_summary_input(source, MAX_CHARS)

    assert prepared.skipped is not None
    assert prepared.skipped.reason is SkipReason.EXTRACTION_FAILED
