"""보안 불변식 — 실행 중 `--ollama-url`(localhost) 외 네트워크 연결 0 (FR-019 / 스펙 §3-6, §4.5).

실제 소켓 `connect`를 감시(패치)해 연결 목적지를 수집하되 실제로는 원격에 나가지 않게 막는다.
파이프라인 전 구간에서 localhost 외 목적지로의 연결 시도가 없음을 강제한다.
"""

from __future__ import annotations

import json
import socket
from pathlib import Path
from typing import Any
from urllib.parse import urlsplit

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Emu

from corpbrain.core import gateway
from corpbrain.core.config import DEFAULT_EMBED_MODEL, DEFAULT_MODEL, ScanConfig
from corpbrain.core.llm.embed import EmbeddingError, embed
from corpbrain.core.llm.ollama_client import OllamaNotAvailableError, detect
from corpbrain.core.pipeline import run_scan
from corpbrain.core.plan import plan_scan

FIXTURE_CORPUS = Path(__file__).resolve().parents[1] / "fixtures" / "sample_corpus"
LOCALHOST_HOSTS = frozenset({"127.0.0.1", "localhost", "::1"})

SUMMARY_JSON = {
    "title": "제목",
    "one_line_summary": "한 줄",
    "key_points": ["a", "b", "c"],
    "summary": "요약",
    "tags": ["t"],
}


class SocketWatcher:
    """소켓 연결 목적지를 수집하는 감시 장치."""

    def __init__(self) -> None:
        self.addresses: list[tuple[str, int]] = []

    def record(self, address: Any) -> None:
        if isinstance(address, tuple) and len(address) >= 2:
            self.addresses.append((str(address[0]), int(address[1])))

    def offenders(self, allowed_hosts: frozenset[str]) -> list[tuple[str, int]]:
        return [addr for addr in self.addresses if addr[0] not in allowed_hosts]


@pytest.fixture
def watch_sockets(monkeypatch: pytest.MonkeyPatch) -> SocketWatcher:
    watcher = SocketWatcher()

    def _fake_connect(self: socket.socket, address: Any) -> None:
        watcher.record(address)
        raise ConnectionRefusedError("보안 감시: 실제 연결을 차단했습니다")

    def _fake_create_connection(address: Any, *args: Any, **kwargs: Any) -> Any:
        watcher.record(address)
        raise ConnectionRefusedError("보안 감시: 실제 연결을 차단했습니다")

    monkeypatch.setattr(socket.socket, "connect", _fake_connect)
    monkeypatch.setattr(socket, "create_connection", _fake_create_connection)
    return watcher


def _host_of(url: str) -> str:
    return urlsplit(url).hostname or ""


def test_pipeline_makes_no_non_localhost_connections(
    watch_sockets: SocketWatcher, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """관문을 스텁한 정상 실행(요약+임베딩 포함)에서 localhost 외 연결 시도가 0건이다.

    v0.4 스펙 §3 완료의 정의 9: 임베딩 API 호출을 포함해 모든 네트워크 호출이 단일 관문
    (`gateway.request_json`)을 경유하는지 검증한다 — 임베딩 엔드포인트가 실제로 호출됐음을
    `calls`로 확인해, 이 테스트가 임베딩 경로를 우회하지 않고 있음을 보장한다.
    """
    calls: list[str] = []

    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        calls.append(url)
        if url.endswith("/api/tags"):
            return {"models": [{"name": DEFAULT_MODEL}, {"name": DEFAULT_EMBED_MODEL}]}
        if url.endswith("/api/embeddings"):
            return {"embedding": [0.1, 0.2, 0.3]}
        return {"response": json.dumps(SUMMARY_JSON, ensure_ascii=False)}

    monkeypatch.setattr(gateway, "request_json", _request_json)

    result = run_scan(
        ScanConfig(folder=FIXTURE_CORPUS, out_dir=tmp_path / "wiki", force_gates=True)
    )

    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []
    assert result.embedding_failures == []
    assert any(url.endswith("/api/embeddings") for url in calls)  # 임베딩 경로가 실제로 돎


def _office_corpus(root: Path) -> None:
    """`.xlsx`·`.xlsm`·`.pptx`를 인라인 생성한다 (v0.8 §3 — 픽스처 폴더를 확장하지 않는다).

    두 오피스 추출기는 zip을 열고 XML을 파싱한다. 그 라이브러리들이 원격 스키마·DTD를 물어
    오지 않는지는 코드를 읽어서가 아니라 **실행 중 소켓을 감시해** 확인해야 하는 종류의
    성질이다 — 새 파싱 의존성을 둘 들이는 이번 슬라이스에서 이 케이스가 필요한 이유다.
    """
    workbook = Workbook()
    workbook.active.append(["항목", "금액"])
    workbook.save(str(root / "예산.xlsx"))
    workbook.save(str(root / "매크로.xlsm"))

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(3_000_000), Emu(500_000))
    box.text_frame.text = "분기 실적"
    presentation.save(str(root / "발표.pptx"))


def test_office_formats_make_no_non_localhost_connections(
    watch_sockets: SocketWatcher, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """v0.8 §3 항목10 — 신규 3종을 포함한 코퍼스의 `scan`에서 소켓 목적지가 localhost뿐이다.

    새 파일로 분리하지 않는다 — 소켓 패치 관용구가 복제되고, 감시장치가 공허하게 통과하지
    않음을 증명하는 `test_watcher_flags_a_gateway_bypass`의 보호를 받지 못한다
    (v0.6·v0.7 결정 계승).

    세 파일이 실제로 요약까지 갔음을 `calls`로 확인해, 이 테스트가 추출 단계에서 전부
    스킵된 채 공허하게 통과하지 않도록 한다.
    """
    calls: list[str] = []

    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        calls.append(url)
        if url.endswith("/api/tags"):
            return {"models": [{"name": DEFAULT_MODEL}, {"name": DEFAULT_EMBED_MODEL}]}
        if url.endswith("/api/embeddings"):
            return {"embedding": [0.1, 0.2, 0.3]}
        return {"response": json.dumps(SUMMARY_JSON, ensure_ascii=False)}

    monkeypatch.setattr(gateway, "request_json", _request_json)
    corpus = tmp_path / "corpus"
    corpus.mkdir()
    _office_corpus(corpus)

    result = run_scan(
        ScanConfig(folder=corpus, out_dir=tmp_path / "wiki", force_gates=True)
    )

    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []
    assert result.skipped == []
    assert len(result.generated) == 3  # 오피스 3종이 실제로 요약까지 갔다


def test_graph_stage_makes_no_non_localhost_connections(
    watch_sockets: SocketWatcher, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """v0.6 §3 항목8 — 그래프 빌드·「관련 문서」 주입 단계에서도 localhost 외 연결이 0건이다.

    그래프 계층은 네트워크 라이브러리를 직접 호출하지 않는다 (v0.6 §4.8). 그래프가 실제로
    만들어졌음을 함께 단언해, 이 테스트가 그래프 경로를 지나치지 않고 있음을 보장한다 —
    감시장치가 공허하게 통과하는 것을 막는 `test_watcher_flags_a_gateway_bypass`와 같은 취지다.
    """

    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        if url.endswith("/api/tags"):
            return {"models": [{"name": DEFAULT_MODEL}, {"name": DEFAULT_EMBED_MODEL}]}
        if url.endswith("/api/embeddings"):
            return {"embedding": [0.1, 0.2, 0.3]}
        return {"response": json.dumps(SUMMARY_JSON, ensure_ascii=False)}

    monkeypatch.setattr(gateway, "request_json", _request_json)

    result = run_scan(
        ScanConfig(folder=FIXTURE_CORPUS, out_dir=tmp_path / "wiki", force_gates=True)
    )

    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []
    assert result.graph is not None
    assert result.graph.stats is not None
    assert result.graph.stats.edges > 0  # 그래프 단계가 실제로 돌았다
    assert result.graph.injection_failures == []


def test_graph_command_opens_no_sockets_at_all(
    watch_sockets: SocketWatcher, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """v0.6 §3 항목8 — `corpbrain graph` 는 소켓을 **하나도** 열지 않는다.

    §4.7에서 순수 조회로 확정했으므로 네트워크가 아예 없어야 한다. `plan_scan`에 대해 이미
    세워 둔 "소켓 0건" 단언과 같은 종류이며, 그래프가 실제로 조회됐음을 함께 단언해 이
    테스트가 조회 경로를 지나치지 않고 있음을 보장한다.
    """
    from corpbrain import cli
    from corpbrain.core.graphstore import SqliteGraphStore, graph_path_for
    from corpbrain.core.models import DocFacts, EdgeType, GraphEdge, GraphNode, NodeType

    out_dir = tmp_path / "wiki"
    out_dir.mkdir()
    source = "/work/docs/a.md"
    with SqliteGraphStore(graph_path_for(out_dir)) as store:
        store.upsert_facts(DocFacts(doc_id=source, title="A", tags=["t"]))
        store.replace_graph(
            [
                GraphNode(id=source, type=NodeType.DOCUMENT, label="A"),
                GraphNode(id="tag:t", type=NodeType.TAG, label="t"),
            ],
            [GraphEdge(src=source, dst="tag:t", type=EdgeType.TAGGED_WITH)],
        )

    def _forbidden(*args: Any, **kwargs: Any) -> Any:
        raise AssertionError("graph 조회는 관문을 부르지 않는다")

    monkeypatch.setattr(gateway, "request_json", _forbidden)

    for view in (["--stats"], ["--central"], ["--neighbors", source]):
        assert cli.main(["graph", "--out", str(out_dir), *view]) == 0

    assert watch_sockets.addresses == []  # localhost조차 열지 않는다


def _seed_search_corpus(out_dir: Path) -> None:
    """벡터 인덱스 + 그래프 DB를 심는다 — 확산이 실제로 일어나는 최소 코퍼스다."""
    from corpbrain.core.graphstore import SqliteGraphStore, graph_path_for
    from corpbrain.core.models import EdgeType, GraphEdge, GraphNode, NodeType
    from corpbrain.core.vectorstore import SqliteVectorStore, index_path_for

    store = SqliteVectorStore(index_path_for(out_dir))
    store.set_model_name(DEFAULT_EMBED_MODEL)
    store.upsert("/docs/a.md", [1.0, 0.0], {"title": "A", "source_path": "/docs/a.md"})
    store.upsert("/docs/b.md", [0.0, 1.0], {"title": "B", "source_path": "/docs/b.md"})
    store.close()

    with SqliteGraphStore(graph_path_for(out_dir)) as graph_store:
        graph_store.replace_graph(
            [
                GraphNode(id="/docs/a.md", type=NodeType.DOCUMENT, label="A"),
                GraphNode(id="/docs/b.md", type=NodeType.DOCUMENT, label="B"),
                GraphNode(id="/docs/c.md", type=NodeType.DOCUMENT, label="C"),
                GraphNode(id="tag:t", type=NodeType.TAG, label="t"),
            ],
            [
                GraphEdge(src="/docs/a.md", dst="tag:t", type=EdgeType.TAGGED_WITH),
                GraphEdge(src="/docs/c.md", dst="tag:t", type=EdgeType.TAGGED_WITH),
            ],
        )


def test_search_connects_once_for_the_query_embedding_and_the_graph_opens_nothing(
    watch_sockets: SocketWatcher, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """v0.7 §3 항목11 — `search` 1회의 관문 호출은 쿼리 임베딩 **1회뿐**이다.

    그래프 확산은 저장된 노드·엣지를 읽기만 하므로 추가 LLM 호출도 추가 네트워크 연결도
    만들지 않는다. 확산이 실제로 일어났음을 함께 단언해 이 불변식이 공허하게 통과하지
    않게 한다 — `test_watcher_flags_a_gateway_bypass`와 같은 취지다.
    """
    from corpbrain.core.search import search_index

    calls: list[str] = []

    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        calls.append(url)
        return {"embedding": [1.0, 0.0]}

    monkeypatch.setattr(gateway, "request_json", _request_json)

    out_dir = tmp_path / "wiki"
    _seed_search_corpus(out_dir)
    results = search_index(out_dir, "질의", top_k=2, graph_decay=0.7)

    assert len(calls) == 1  # 쿼리 임베딩 1회 — 그래프 단계는 관문을 부르지 않는다
    assert calls[0].endswith("/api/embeddings")
    assert any(result.expansion is not None for result in results)  # 확산이 실제로 돌았다
    assert watch_sockets.addresses == []  # 관문을 우회해 직접 연 소켓이 없다


def test_search_only_dials_the_given_localhost_url(
    watch_sockets: SocketWatcher, tmp_path: Path
) -> None:
    """v0.7 §3 항목11 — 관문 스텁 없이 `search`가 실제로 여는 소켓은 localhost뿐이다."""
    from corpbrain.core.errors import PreconditionError
    from corpbrain.core.search import search_index

    out_dir = tmp_path / "wiki"
    _seed_search_corpus(out_dir)

    with pytest.raises(PreconditionError):
        search_index(out_dir, "질의", top_k=2, graph_decay=0.7)

    assert watch_sockets.addresses  # 실제로 연결을 시도했다
    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []
    assert {addr[1] for addr in watch_sockets.addresses} == {11434}


def test_gateway_only_dials_the_given_localhost_url(watch_sockets: SocketWatcher) -> None:
    """관문 스텁 없이 detect가 실제로 여는 소켓은 지정한 localhost 주소뿐이다."""
    url = "http://127.0.0.1:11434"

    with pytest.raises(OllamaNotAvailableError):
        detect(url)

    assert watch_sockets.addresses  # 실제로 연결을 시도했다
    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []
    assert all(host == _host_of(url) for host, _ in watch_sockets.addresses)


def test_embed_only_dials_the_given_localhost_url(watch_sockets: SocketWatcher) -> None:
    """관문 스텁 없이 embed()가 실제로 여는 소켓은 지정한 localhost 주소뿐이다 (v0.4 §3 항목9)."""
    url = "http://127.0.0.1:11434"

    with pytest.raises(EmbeddingError):
        embed("문서 요약 텍스트", DEFAULT_EMBED_MODEL, url)

    assert watch_sockets.addresses  # 실제로 연결을 시도했다
    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []
    assert all(host == _host_of(url) for host, _ in watch_sockets.addresses)


def test_custom_localhost_port_is_the_only_target(watch_sockets: SocketWatcher) -> None:
    """`--ollama-url`을 다른 localhost 포트로 바꿔도 그 주소로만 연결한다(외부 호스트 하드코딩 부재)."""
    url = "http://127.0.0.1:9999"

    with pytest.raises(OllamaNotAvailableError):
        detect(url)

    assert {addr[1] for addr in watch_sockets.addresses} == {9999}
    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []


def test_watcher_flags_a_gateway_bypass(watch_sockets: SocketWatcher) -> None:
    """관문을 우회해 외부로 직접 연결하면 감시 장치가 검출한다(회귀 방지 자기검증)."""
    with pytest.raises(OSError):
        socket.create_connection(("example.com", 443))

    assert watch_sockets.offenders(LOCALHOST_HOSTS) == [("example.com", 443)]


def test_plan_scan_opens_no_sockets_and_bypasses_gateway(
    watch_sockets: SocketWatcher, tmp_path: Path
) -> None:
    """v0.2 완료의 정의 3: plan은 localhost 포함 어떤 소켓도 열지 않고 관문도 거치지 않는다.

    하드웨어 감지는 로컬 nvidia-smi subprocess만 쓰므로(소켓 아님) 이 실행에서 소켓 연결이
    전혀 없어야 하고, `gateway.requested_urls()`도 비어 있어야 한다.
    """
    gateway.reset_requested_urls()
    (tmp_path / "note.txt").write_text("본문", encoding="utf-8")
    (tmp_path / "report.pdf").write_bytes(b"%PDF-1.4 stub")

    plan_scan(ScanConfig(folder=tmp_path))

    assert watch_sockets.addresses == []  # localhost 포함 0건
    assert gateway.requested_urls() == ()


# --- v0.5: 엔진별 목적지 불변식 (스펙 §3 항목7) --------------------------------

ANTHROPIC_HOSTS = frozenset({"api.anthropic.com"})


def test_cloud_run_reaches_only_the_allowlisted_host(
    watch_sockets: SocketWatcher, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """v0.5 항목7: `--engine cloud` 실행 중 목적지는 localhost(임베딩)와 allowlist 호스트뿐이다.

    임베딩은 엔진과 무관하게 로컬이므로(§2 비목표) localhost도 정상 목적지다. 여기서 잡고자
    하는 것은 "그 둘 외의 어떤 목적지도 없다"는 불변식이다.
    """
    from corpbrain.core.config import ENGINE_CLOUD
    from corpbrain.core.consent import grant_cloud_consent
    from corpbrain.core.llm.anthropic_client import API_KEY_ENV_VAR, SUMMARY_TOOL_NAME

    home = tmp_path / "home"
    home.mkdir()
    monkeypatch.setattr(Path, "home", classmethod(lambda _cls: home))
    monkeypatch.setenv(API_KEY_ENV_VAR, "sk-test")
    grant_cloud_consent()

    calls: list[str] = []

    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        calls.append(url)
        if url.endswith("/api/tags"):
            return {"models": [{"name": DEFAULT_MODEL}, {"name": DEFAULT_EMBED_MODEL}]}
        if url.endswith("/api/embeddings"):
            return {"embedding": [0.1, 0.2, 0.3]}
        if url.endswith("/v1/models"):
            return {"data": []}
        return {
            "stop_reason": "tool_use",
            "content": [
                {"type": "tool_use", "name": SUMMARY_TOOL_NAME, "input": SUMMARY_JSON}
            ],
        }

    monkeypatch.setattr(gateway, "request_json", _request_json)

    run_scan(
        ScanConfig(
            folder=FIXTURE_CORPUS,
            out_dir=tmp_path / "wiki",
            engine=ENGINE_CLOUD,
            force_gates=True,
        )
    )

    assert watch_sockets.offenders(LOCALHOST_HOSTS | ANTHROPIC_HOSTS) == []
    # 클라우드 경로가 실제로 돌았음을 확인 — 불변식이 공허하게 통과하지 않게 한다.
    assert any(url.endswith("/v1/messages") for url in calls)
    assert all(_host_of(url) in LOCALHOST_HOSTS | ANTHROPIC_HOSTS for url in calls)


def test_local_run_never_reaches_the_cloud_host(
    watch_sockets: SocketWatcher, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """v0.5 항목7: 기본(`--engine` 미지정) 실행은 v0.4까지와 동일하게 localhost 외 연결이 없다."""
    calls: list[str] = []

    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        calls.append(url)
        if url.endswith("/api/tags"):
            return {"models": [{"name": DEFAULT_MODEL}, {"name": DEFAULT_EMBED_MODEL}]}
        if url.endswith("/api/embeddings"):
            return {"embedding": [0.1, 0.2, 0.3]}
        return {"response": json.dumps(SUMMARY_JSON, ensure_ascii=False)}

    monkeypatch.setattr(gateway, "request_json", _request_json)

    run_scan(
        ScanConfig(folder=FIXTURE_CORPUS, out_dir=tmp_path / "wiki", force_gates=True)
    )

    assert watch_sockets.offenders(LOCALHOST_HOSTS) == []
    assert all(_host_of(url) in LOCALHOST_HOSTS for url in calls)


def test_gateway_blocks_a_cloud_call_to_a_wrong_host(watch_sockets: SocketWatcher) -> None:
    """NetworkGuard가 allowlist 밖 목적지를 소켓 이전에 막는다 (자기검증)."""
    with pytest.raises(gateway.NetworkGuardError):
        gateway.request_json(
            "https://evil.example.com/v1/messages",
            allowed_hosts=("api.anthropic.com",),
            require_https=True,
        )

    assert watch_sockets.addresses == []  # 소켓을 아예 열지 않았다
