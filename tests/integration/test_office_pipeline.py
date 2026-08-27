"""통합 테스트 — v0.8 `.xlsx`/`.xlsm`/`.pptx`가 파이프라인을 타는지 코어 API 직접 호출로 검증.

완료의 정의 1(출력 트리)·3(수식 전용 통합문서의 `empty_document`)·5(암호화 detail과 부분
성공)·9(`max_chars` 상한)를 덮는다.

외부호출 단일 관문(`gateway.request_json`)만 스텁하고 `run_scan`을 직접 부른다 (스펙 §4.5).
픽스처는 `tmp_path`에 `openpyxl`·`python-pptx`로 인라인 생성한다 — `tests/fixtures/sample_corpus/`를
확장하지 않는 것은 세 테스트 파일이 그 폴더를 공유하고 있고, 2026-08-22에 스모크 산출물이
그 폴더를 오염시켜 통합테스트 3건을 깨뜨린 사고가 있었기 때문이다 (스펙 §3).

픽스처 생성 헬퍼가 `tests/unit/test_extract.py`와 복제된 것은 의도적이다 — `.pdf`가 이미
같은 선택을 했고, 공유 인프라를 새로 들이면 `.pdf`는 복제, 오피스는 공유가 되어 관용구가
둘로 갈린다 (스펙 §3 「헬퍼 배치」).
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Emu

from corpbrain.core import gateway, run_scan
from corpbrain.core.config import DEFAULT_EMBED_MODEL, DEFAULT_MODEL, ScanConfig
from corpbrain.core.render import FRONT_MATTER_KEYS, SECTION_HEADERS

#: OLE 복합문서(CFBF) 매직 넘버 — 암호화된 OOXML이 zip 대신 이 컨테이너로 감싸진다.
#: **실제 암호화 통합문서를 만들 필요가 없다** (스펙 §3 항목5) — 판정 근거가 이 8바이트뿐이다.
OLE_SIGNATURE_BYTES = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

TAGS_RESPONSE = {"models": [{"name": DEFAULT_MODEL}, {"name": DEFAULT_EMBED_MODEL}]}

SUMMARY_JSON = {
    "title": "오피스 통합 제목",
    "one_line_summary": "한 줄 요약.",
    "key_points": ["p1", "p2", "p3"],
    "summary": "문단 요약.",
    "tags": ["t1", "t2"],
    "entities": ["인사팀"],
}


def _write_xlsx(path: Path, sheets: dict[str, list[list[object]]]) -> None:
    workbook = Workbook()
    workbook.remove(workbook.active)
    for title, rows in sheets.items():
        worksheet = workbook.create_sheet(title)
        for row in rows:
            worksheet.append(row)
    workbook.save(str(path))


def _write_formula_only_xlsx(path: Path) -> None:
    """수식만 있고 계산값 캐시가 없는 통합문서 — Excel이 저장하지 않은 파일의 재현이다."""
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "계산"
    worksheet["A1"] = "=SUM(B1:B9)"
    worksheet["A2"] = "=A1*2"
    workbook.save(str(path))


def _write_pptx(path: Path, slides: list[list[str]]) -> None:
    presentation = Presentation()
    for texts in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for index, text in enumerate(texts):
            box = slide.shapes.add_textbox(
                Emu(0), Emu(index * 1_000_000), Emu(3_000_000), Emu(500_000)
            )
            box.text_frame.text = text
    presentation.save(str(path))


def _ok_gateway(monkeypatch: pytest.MonkeyPatch) -> None:
    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        if url.endswith("/api/tags"):
            return TAGS_RESPONSE
        if url.endswith("/api/embeddings"):
            return {"embedding": [0.1, 0.2, 0.3]}
        return {"response": json.dumps(SUMMARY_JSON, ensure_ascii=False)}

    monkeypatch.setattr(gateway, "request_json", _request_json)


def test_office_formats_generate_the_expected_output_tree(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """DoD 1 — 신규 3종에 위키가 1개씩 생기고 출력 트리가 기대 집합과 정확히 일치한다.

    기존 지원 4종을 함께 넣어 **하위 호환**도 같은 실행에서 확인한다 — 신규 확장자 추가가
    기존 포맷의 동작을 바꾸지 않아야 한다.
    """
    _ok_gateway(monkeypatch)
    corpus = tmp_path / "corpus"
    (corpus / "sub").mkdir(parents=True)
    _write_xlsx(corpus / "예산.xlsx", {"인사": [["항목", "금액"], ["인건비", 120000.0]]})
    _write_xlsx(corpus / "sub" / "매크로.xlsm", {"시트": [["매크로 통합문서"]]})
    _write_pptx(corpus / "발표.pptx", [["분기 실적"], ["다음 계획"]])
    (corpus / "메모.txt").write_text("평문 메모", encoding="utf-8")
    (corpus / "readme.md").write_text("# 마크다운", encoding="utf-8")
    out_dir = tmp_path / "wiki"

    result = run_scan(ScanConfig(folder=corpus, out_dir=out_dir, force_gates=True))

    produced = {
        path.relative_to(out_dir).as_posix()
        for path in out_dir.rglob("*.md")
    }
    assert produced == {
        "예산.xlsx.md",
        "sub/매크로.xlsm.md",
        "발표.pptx.md",
        "메모.txt.md",
        "readme.md.md",
    }
    assert result.skipped == []

    for name in ("예산.xlsx.md", "발표.pptx.md"):
        markdown = (out_dir / name).read_text(encoding="utf-8")
        for key in FRONT_MATTER_KEYS:
            assert f"{key}:" in markdown
        for header in SECTION_HEADERS:
            assert header in markdown


def test_formula_only_workbook_is_skipped_as_empty_document(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """DoD 3 — 수식만 있고 캐시된 계산값이 없는 `.xlsx`는 `empty_document`로 스킵된다.

    이 항목이 성립하는 것은 §4.2 T1(「내용 없는 시트에는 경계 줄도 내지 않는다」) 덕분이다.
    경계 줄을 늘 냈다면 추출 결과가 `[시트: 계산]` 한 줄로 남아 `text.strip()` 검사를
    통과하고, 시트 이름만 든 입력이 LLM까지 갔을 것이다.
    """
    _ok_gateway(monkeypatch)
    corpus = tmp_path / "corpus"
    corpus.mkdir()
    _write_formula_only_xlsx(corpus / "수식뿐.xlsx")
    out_dir = tmp_path / "wiki"

    result = run_scan(ScanConfig(folder=corpus, out_dir=out_dir, force_gates=True))

    assert [skip.reason.value for skip in result.skipped] == ["empty_document"]
    assert not (out_dir / "수식뿐.xlsx.md").exists()
    assert result.generated == []


def test_ole_and_corrupted_office_files_get_different_details(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """DoD 5 — 암호화(OLE)와 손상이 서로 다른 detail을 받고, 나머지 파일은 정상 처리된다.

    개별 파일 실패를 전체 실패로 위장하지 않는다 (v0.1 §5 계승).
    """
    _ok_gateway(monkeypatch)
    corpus = tmp_path / "corpus"
    corpus.mkdir()
    (corpus / "잠긴통합문서.xlsx").write_bytes(OLE_SIGNATURE_BYTES + b"\x00" * 64)
    (corpus / "잠긴발표.pptx").write_bytes(OLE_SIGNATURE_BYTES + b"\x00" * 64)
    (corpus / "손상.xlsx").write_bytes(b"PK\x03\x04 not really a zip")
    _write_xlsx(corpus / "정상.xlsx", {"시트": [["정상 내용"]]})
    out_dir = tmp_path / "wiki"

    result = run_scan(ScanConfig(folder=corpus, out_dir=out_dir, force_gates=True))

    skipped = {skip.path.name: skip for skip in result.skipped}
    assert set(skipped) == {"잠긴통합문서.xlsx", "잠긴발표.pptx", "손상.xlsx"}
    for name in ("잠긴통합문서.xlsx", "잠긴발표.pptx"):
        assert skipped[name].reason.value == "extraction_failed"
        assert "암호화되었거나 구형 이진 포맷" in skipped[name].detail
    assert skipped["손상.xlsx"].reason.value == "extraction_failed"
    assert "암호화되었거나 구형 이진 포맷" not in skipped["손상.xlsx"].detail

    # 부분 성공 — 나머지 파일은 그대로 처리된다.
    assert (out_dir / "정상.xlsx.md").exists()
    assert {wiki.source_path.name for wiki in result.generated} == {"정상.xlsx"}


def test_office_extraction_respects_max_chars(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """DoD 9 — `max_chars`를 초과하는 `.xlsx`/`.pptx`의 요약 입력 길이가 상한 이하다.

    관문 스텁이 실제로 넘어온 프롬프트를 붙잡아 길이를 잰다 — 추출기 단위테스트가 아니라
    **파이프라인이 실제로 넘기는 값**을 본다.
    """
    max_chars = 400
    prompts: list[str] = []

    def _request_json(url: str, *, method: str = "GET", payload: Any = None, **_: Any) -> Any:
        if url.endswith("/api/tags"):
            return TAGS_RESPONSE
        if url.endswith("/api/embeddings"):
            return {"embedding": [0.1, 0.2, 0.3]}
        prompts.append(str((payload or {}).get("prompt", "")))
        return {"response": json.dumps(SUMMARY_JSON, ensure_ascii=False)}

    monkeypatch.setattr(gateway, "request_json", _request_json)

    corpus = tmp_path / "corpus"
    corpus.mkdir()
    # 프롬프트 템플릿(한국어)에 섞이지 않는 마커 글자를 쓴다 — 그래야 프롬프트에서 본문
    # 글자수만 정확히 셀 수 있다. 두 포맷을 다른 마커로 갈라 서로의 분량이 섞이지 않게 한다.
    _write_xlsx(corpus / "큰통합문서.xlsx", {"시트": [["X" * 200] for _ in range(40)]})
    _write_pptx(corpus / "큰발표.pptx", [["Y" * 200] for _ in range(40)])
    out_dir = tmp_path / "wiki"

    result = run_scan(
        ScanConfig(
            folder=corpus, out_dir=out_dir, max_chars=max_chars, force_gates=True
        )
    )

    assert len(result.generated) == 2
    assert len(prompts) == 2
    marker_counts = sorted(
        max(prompt.count("X"), prompt.count("Y")) for prompt in prompts
    )
    # 두 문서 모두 본문이 8,000자였으나 요약 입력에는 상한까지만 실렸다.
    assert all(count <= max_chars for count in marker_counts)
    # 상한 자체가 무의미해지지 않았음을 함께 본다 — 절단이 실제로 일어났다.
    assert all(count > 0 for count in marker_counts)
    assert sum(prompt.count("X") for prompt in prompts) < 40 * 200
